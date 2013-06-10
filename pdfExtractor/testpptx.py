from pptx import Presentation

prs = Presentation("/home/akhsig/code/pdfExtractor/files/2_Chemical Signalling_3204_Winter 2012_CuLearn.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_textframe:
            continue
        for paragraph in shape.textframe.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
print text_runs
