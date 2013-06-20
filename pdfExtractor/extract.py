###################################################################
##                  DocumentContentExtractor                     ##
##  Author: Olivier Gilbert, contact: ogilbert@e-medhosting.com  ##
##              Date: May 16th, 2013.  Version: 1.1              ##
###################################################################

###################################################################
##                           Changelog                           ##
## - v0.1:                                                       ##
##     	=> script can read pdf files                             ##
## - v0.2                                                        ##
##	=> script attempts to decrypt pdf files                  ##
## - v1.0                                                        ##
## 	=> script now reads ppt, pptx, doc, docx                 ##
## - v1.1                                                        ##
##	=> script now shows progress bar                         ##
## - v2.0							 ##
##	=> script now takes full csv file and tries to find	 ##
##	   information on presentation based on name		 ##
###################################################################

# This script takes a list of pdf files (will us ls of a given directory for list of files)
# and extracts the text content of each pdf file and writes it on a new line in an output file
# If the pdf is encrypted, the script will attempt to decrypt it with the empty password
# (since we don't need to edit, this should work for most documents!)

# TOOLS NEEDED
# pypdf2: https://github.com/knowah/PyPDF2/
# catppt suite: http://www.wagner.pp.ru/~vitus/software/catdoc/
# python-pptx: https://python-pptx.readthedocs.org/en/latest/index.html
# python-mysql: http://sourceforge.net/projects/mysql-python/

# First import required packages
import PyPDF2 as p #must be installed for this script to work
from subprocess import check_output # needed to run shell commands
import subprocess as sp #need this for the subprocess module in general
import string as s
from pptx import Presentation #needed for pptx support!
import sys #seems to be needed for fixing problem with the try statement for pptx
import os #will be used for directory listing
import MySQLdb #needed for mysql integration

host = "lecaire"
port = 3306
user = "ogilbert"
password = "thrhhhvthdb"
schema = "multiwebcast"

db = db=MySQLdb.connect(host=host, user=user, passwd=password, port=port, db=schema);
cursor = db.cursor()

# now that we have the db integration, we create a dict with id from csv file and tuple with all the information
files = []
texts = {}
problems = []
toolbar_width = 100
extract_dir = "/home/akhsig/git/code-home/pdfExtractor/files/SIU" # variable for the absolute path to directory for pdf files!
output_files = open('/home/akhsig/git/code-home/pdfExtractor/output_files.txt', 'w')
output_texts = open('/home/akhsig/git/code-home/pdfExtractor/output_texts.txt', 'w')
output_problems = open('/home/akhsig/git/code-home/pdfExtractor/output_problems.txt', 'w')

#get a list of the words on the GSL
#this reduces the number of elements 
#in the set of words from each document
#because we assume these words are less
#important; to be discussed?
gsl = open('gsl', 'r')
gsl_words_list = []
for line in gsl:
	gsl_words_list.append(line.split()[2])
gsl.close()

#now make it a set!
gsl_words = set(gsl_words_list)

prefix = extract_dir
if prefix[len(prefix)-1] != "/":
	prefix += "/"

pres_dict = {}
presentations = []
# first thing we do is create a basic dict with id and in tuple: title, location

path_to_csv = "/home/akhsig/git/code-home/pdfExtractor/indexation.csv"
csv = open(path_to_csv, 'r')
for line in csv:
	fields = s.split(line, ',')
	docid = fields[0]
	title = fields[1]
	loc = fields[2]
	location = prefix+s.split(loc, '/')[-1]
	pres_dict[title] = location
	presentations.append(title.strip('"'))

in_filter = '"'+'","'.join(presentations)+'"'

query = """     SELECT
                        _40_conference.c_id,
			_40_conference.c_name
                FROM
                        _40_conference
		INNER JOIN
			_40_conference_event_group
			ON _40_conference.ce_id = _40_conference_event_group.ce_id
                WHERE
			_40_conference_event_group.g_id = 184 AND
                        _40_conference.c_name in ("""
query += in_filter
query += ')'
cursor.execute(query)

result = cursor.fetchall()
result_dict = {x[1]: x[0] for x in result}

# in result_dict, we now have all the information for each presentation indexed by presentation title (found in presentations)

print "Reading and parsing files"
 
files_len = float(len(presentations))
current = 0
tmp = 0

for pres in presentations:
	if pres in pres_dict.keys():
		f = pres_dict[pres]
	else:
		continue
	tmp += 1
	#sys.stdout.write("\b" * (len(str(current))+1))
	text = ""
	ext = f.split('.')[f.count('.')]
	if ext == "ppt":
		#it's a ppt, need to use catppt tool!
		#make sure to handle errors!
		out = check_output(["catppt", f], stderr=sp.STDOUT)
		if out == f+" is not OLE file or Error":
			print f+" will be added to problematic files!"
			problems.append(f)
			#files.remove(f)
		else:
			#we don't have an error, out contains content of document
			out = set(s.replace(out,"\n"," ").replace(".", "").replace(",", "").replace("!", "").replace("?", "").replace(":", "").replace("'", "").encode('utf-8').lower().split())
			#out -= gsl_words
			texts[pres] = s.join(out," ")
	elif ext == "pptx":
		#for pptx, use a try since there seems to be many things that can go wrong with py-pptx...
		try:
			prs = Presentation(f)
		except (LookupError, KeyError):
			#something went wrong along the way, add file to problems and move on
			problems.append(f)
			#files.remove(f)
			sys.exc_clear()
			continue
		#we were able to open the file, read the content!
		#else:
		out = ""
		for slide in prs.slides:
			for shape in slide.shapes:
		        	if not shape.has_textframe:
		            		continue
       				for paragraph in shape.textframe.paragraphs:
					for run in paragraph.runs:
       	        				out += run.text+" "
		out = set(s.replace(out, "\n", " ").replace(".", "").replace(",", "").replace("!", "").replace("?", "").replace(":", "").replace("'", "").encode('utf-8').lower().split())
		#out -= gsl_words
		#texts.append(s.join(out, " "))
		texts[pres] = s.join(out," ")
	elif ext == "doc":
		#we have a doc file, use catdoc! if it doesnt work, we should get a segfault!
		try:
			out = check_output(['catdoc', f], stderr=sp.STDOUT)
		except sp.CalledProcessError:
			problems.append(f)
			sys.exc_clear()
			continue

		if out == "Segmentation fault (core dumped)":
			print f+" will be added to problems!"
			problems.append(f)
			#files.remove(f)
		else:
			#we have an output! parse it a little and add it to texts!
			out = set(s.replace(out, "\n", " ").replace(".", "").replace(",", "").replace("!", "").replace("?", "").replace(":", "").replace("'", "").encode('utf-8').lower().split())
			#out -= gsl_words
			#texts.append(s.join(out, " "))
			texts[pres] = s.join(out," ")
	elif ext == "docx":
		#we have a docx, use the perl script to try and extract content
		out = check_output(['/usr/bin/docx2txt.pl', f, '-'], stderr=sp.STDOUT)
		if out.startswith("Can't read docx file <") or out.startswith("Failed to extract required information from <"):
			#we have an error, add to problems!
			print f+" will be added to problems!"
                        problems.append(f)
			#files.remove(f)http://sourceforge.net/projects/mysql-python/
		else:
			#out contains the output of the file
			out = set(s.replace(out, "\n", " ").replace(".", "").replace(",", "").replace("!", "").replace("?", "").replace(":", "").replace("'", "").encode('utf-8').lower().split())
                        #out -= gsl_words
                        #texts.append(s.join(out, " "))
			texts[pres] = s.join(out," ")
	#elif ext == "xls":  #for now dont need to support these files...
	#elif ext == "xlsx":
	elif ext == "pdf":
		pdf = p.PdfFileReader(open(str(f), "rb"))
		if pdf.isEncrypted:
			# we try and decrypt it!
			pdf.decrypt("")
		
		page_num = pdf.getNumPages()
		for i in range(page_num):
			page = pdf.getPage(i)
			text += s.replace(page.extractText(),"\n"," ").replace(".", "").replace(",", "").replace("!", "").replace("?", "").replace(":", "").replace("'", "")
		text = set(text.encode('utf-8').lower().split())
		#text -= gsl_words
		#texts.append(s.join(text, " "))
		texts[pres] = s.join(text," ")
	else: 
		#file type is something else
		#for now, add file to problems
		problems.append(f)
		#files.remove(f)
	
	temp = tmp/files_len*100
        if temp > current:
		sys.stdout.write(str(int(temp))+"%")
		sys.stdout.flush()
                #for i in range(int(temp-current)):
                #        sys.stdout.write("-")
                #        sys.stdout.flush()
                current = temp

sys.stdout.write("\n")

#now, we should have the content of all the files that could be opened and read and problems should be populated with files that need to be checked manually

print "Updating database!"

'''
print "Printing output to files!"

files = [x for x in files if x not in problems]

for prob in problems:
	output_problems.write(prob+"\n")

for i in range(len(files)):
	output_texts.write(files[i]+":"+texts[i]+"\n")

for f in files:
	output_files.write(f+"\n")

output_problems.close()
output_texts.close()
output_files.close()
'''

insert_query = """

	INSERT INTO `z_siu_content_index_20130614` (
		c_id,
		content
	)
	VALUES
		
"""

for key in result_dict.keys():
	if key in texts.keys():
		insert_query += "("+str(result_dict[key])+",'"+str(texts[key])+"'),\n"

insert_query = insert_query.rstrip(",")

output_texts.write(insert_query)
output_texts.close()

cursor.execute(insert_query)


